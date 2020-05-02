import os
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt, Mm, Inches
import PyPDF2
import re
import xlrd
from PIL import Image
import locale
locale.setlocale(locale.LC_ALL, 'pt_BR.UTF-8')


def formated_number(a):

    num = locale.format_string("%1.2f", float(a), grouping=True)
    return(num)


def apresentacao(agencia,endereco,foto_fachada_png,area_construida,area_terreno,iptu_espelho,mapa_demanda,demanda_descricao,mapa_satelite_png,zoneamento_sim_nao,zon_1,zon_2,zon_3,zon_4,zon_5,zon_6,zon_7,zon_png_1,zon_png_2,zon_png_3,rating,zona_de_influencia_residencial,tabela_residencial_png,zona_de_influencia_comercial,tabela_comercial_png,zona_de_influencia_terreno,tabela_terreno_png,mediana_terreno,valor_terreno_evolutivo,avaliacao_evolutivo_parte_2,cidade_uf_referencia,m2_construcao,valor_construcao,valor_total_imovel,valor_venal_png,n,m2comercial,valor_comparativo,involutivo_png,involutivo_valor,proj_hiptotetico,proj_area_util,proj_area_total,valor_oferta,parecer,finalizado):
    
    # arquivos
    prs = Presentation()
    path = './'
    
    
    for arquivo in os.listdir(path):
        if arquivo.startswith(agencia) and not arquivo.endswith('.pptx'):
            folder = arquivo + '/'
    # print(folder)
    path_slides_padrao = 'slides_padrao/'
    slides_padrao = os.listdir(path_slides_padrao)
    
    for arq in slides_padrao:
        if 'slide_1.png' in arq:
            slide_1 = arq
        if 'slide_2.png' in arq:
            slide_2 = arq
        if 'slide_3.png' in arq:
            slide_3 = arq
        if 'slide_4.png' in arq:
            slide_4 = arq
        if 'slide_5.png' in arq:
            slide_5 = arq
        if 'slide_7.png' in arq:
            slide_7 = arq
        if 'bg_texto.png' in arq:
            bg = arq
        if 'casa' in arq:
            img_casa = arq
        if 'avaliacao_comparativo_formulas_1' in arq:
            avaliacao_comparativo_formulas_1 = arq
        if 'avaliacao_comparativo_formulas_2' in arq:
            avaliacao_comparativo_formulas_2 = arq
        if 'parceiros' in arq:
            parceiros = arq
        if 'texto_avaliacao_evolutivo_parte_2' in arq:
            texto_avaliacao_evolutivo_parte_2 = arq
        if 'calculos_avaliatorios' in arq:
            calculos_avaliatorios = arq
        if 'consideracoes' in arq:
            consideracoes = arq

    num_pag = 1
########################################################
    ##### SLIDE 1 - Intro
    # background
    slide_1_path = path_slides_padrao + slide_1
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    slide_1_left = slide_1_top = Inches(0)
    slide_1 = slide.shapes.add_picture(slide_1_path, slide_1_left, slide_1_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(slide_1._element)
    slide.shapes._spTree.insert(2, slide_1._element)


    # Numeracao da página
    # text_top = Mm(127)
    # text_left = Mm(0)
    # rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    # tf = rodape.text_frame
    # p = tf.add_paragraph()
    # run = p.add_run()
    # run.text = str(num_pag)
    # p.alignment = PP_ALIGN.RIGHT
    # font = run.font
    # font.name = 'Century Gothic'
    # font.size = Pt(12)
    # font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

##########################################################
    ##### SLIDE 2 - Uma startup de Análise
    # background 
    slide_2_path = path_slides_padrao + slide_2
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    slide_2_left = slide_2_top = Inches(0)
    slide_2 = slide.shapes.add_picture(slide_2_path, slide_2_left, slide_2_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(slide_2._element)
    slide.shapes._spTree.insert(2, slide_2._element)

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

#############################################################
    ##### SLIDE 3 - EEMOVEL INTELIGÊNCIA IMOBILIÁRIA (DATA LAKE)
    # background
    slide_3_path = path_slides_padrao + slide_3
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    slide_3_left = slide_3_top = Inches(0)
    slide_3 = slide.shapes.add_picture(slide_3_path, slide_3_left, slide_3_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(slide_3._element)
    slide.shapes._spTree.insert(2, slide_3._element)

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    ##### SLIDE 4 - VERTICAIS
    # background
    slide_4_path = path_slides_padrao + slide_4
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    slide_4_left = slide_4_top = Inches(0)
    slide_4 = slide.shapes.add_picture(slide_4_path, slide_4_left, slide_4_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(slide_4._element)
    slide.shapes._spTree.insert(2, slide_4._element)

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(25,37,54)
    num_pag += 1

###################################################################
    ##### SLIDE 5 - VALUATION INTELLIGENCE
    # background
    slide_5_path = path_slides_padrao + slide_5
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    slide_5_left = slide_5_top = Inches(0)
    slide_5 = slide.shapes.add_picture(slide_5_path, slide_5_left, slide_5_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(slide_5._element)
    slide.shapes._spTree.insert(2, slide_5._element)

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    ##### SLIDE 6 - FACHADA
    # background
    bg_path = path_slides_padrao + bg
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_left = bg_top = Inches(0)
    bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # imagem inserida Fachada
    fachada_path = str(folder) + str(foto_fachada_png)
    img_fachada = Image.open(fachada_path)
    img_fachada_width, img_fachada_height = img_fachada.size
    fachada_top = Mm(10)
    height = Mm(85)
    width = (height * img_fachada_width)/img_fachada_height
    fachada_left = (prs.slide_width - width)/2
    fachada = slide.shapes.add_picture(fachada_path, fachada_left, fachada_top, width=width)

    # INFORMAÇÃO FOTO
    text_top = Mm(90)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'AGÊNCIA ' + str(agencia)
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(26)
    font.color.rgb = RGBColor(25,37,54)
    
    # Endereço
    text_top = Mm(100)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    if len(str(endereco)) > 52:
        font_size = 22
    elif len(str(endereco)) > 60:
        font_size = 20
    else:
        font_size = 26
    run.text = str(endereco)
    
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(font_size)
    font.color.rgb = RGBColor(25,37,54)

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    ##### SLIDE - DATA ANALYSIS
    # background
    slide_7_path = path_slides_padrao + slide_7
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    slide_7_left = slide_7_top = Inches(0)
    slide_7 = slide.shapes.add_picture(slide_7_path, slide_7_left, slide_7_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(slide_7._element)
    slide.shapes._spTree.insert(2, slide_7._element)

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    ##### SLIDE 8
    # background
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_left = bg_top = Inches(0)
    bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # imagem inserida Foto Satélite
    satelite_path = folder + mapa_satelite_png
    img_satelite = Image.open(satelite_path)
    img_satelite_width, img_satelite_height = img_satelite.size
    satelite_top = Mm(18)
    height = Mm(80)
    width = (height * img_satelite_width)/img_satelite_height
    satelite_left = (prs.slide_width - width)/2
    satelite = slide.shapes.add_picture(satelite_path, satelite_left, satelite_top, width=width)

    # Título - Informações do Imóvel
    text_top = Mm(-4)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'INFORMAÇÕES DO IMÓVEL'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)

    # INFORMAÇÃO FOTO - Área do Terreno + Área Construída
    text_top = Mm(93)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    try:
        run.text = 'Área do Terreno: ' + str(formated_number(area_terreno)) + 'm²\nÁrea Construída: ' + str(formated_number(area_construida)) + 'm²'
    except Exception as e:
        print(agencia,'Não encontrou número Área do Terreno ou Área Construída', e)
        pass
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(26)
    font.bold = True
    font.color.rgb = RGBColor(25,37,54)

    

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    ##### SLIDE INFORMAÇÕES DO IMÓVEL
    # background
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_left = bg_top = Inches(0)
    bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # imagem inserida iptu espelho
    area_iptu_path = folder + iptu_espelho
    img_area_iptu = Image.open(area_iptu_path)
    img_area_iptu_width, img_area_iptu_height = img_area_iptu.size
    if img_area_iptu_height > 288:
        height = Mm(90)
        width = (height * img_area_iptu_width)/img_area_iptu_height
    else:
        width = Mm(140)
        height = (width * img_area_iptu_height)/img_area_iptu_width
    
    area_iptu_top = (prs.slide_height - height)/2
    area_iptu_left = (prs.slide_width - width)/2
    area_iptu_left = (prs.slide_width - width)/2
    area_iptu_slide = slide.shapes.add_picture(area_iptu_path, area_iptu_left, area_iptu_top, width=width)

    # Título - Informações do Imóvel
    text_top = Mm(-4)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'INFORMAÇÕES DO IMÓVEL'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    ##### SLIDE INFLUENCIADORES DE DEMANDA - IMAGEM 1
    # background
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_left = bg_top = Inches(0)
    bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # imagem inserida influenciadores de demanda mapa
    mapa_demanda_path = folder + mapa_demanda
    img_mapa_demanda = Image.open(mapa_demanda_path)
    img_mapa_demanda_width, img_mapa_demanda_height = img_mapa_demanda.size
    mapa_demanda_top = Mm(20)
    height = Mm(100)
    width = (height * img_mapa_demanda_width)/img_mapa_demanda_height
    mapa_demanda_left = (prs.slide_width - width)/2
    mapa_demanda_slide = slide.shapes.add_picture(mapa_demanda_path, mapa_demanda_left, mapa_demanda_top, width=width)

    # Título - Informações do Imóvel
    text_top = Mm(-4)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'INFLUENCIADORES DE DEMANDA'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1


###################################################################
    ##### SLIDE INFLUENCIADORES DE DEMANDA - IMAGEM 2
    # background
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_left = bg_top = Inches(0)
    bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # imagem inserida influenciadores de demanda mapa
    demanda_descricao_path = folder + demanda_descricao
    img_demanda_descricao = Image.open(demanda_descricao_path)
    img_demanda_descricao_width, img_demanda_descricao_height = img_demanda_descricao.size
    if img_demanda_descricao_height > 288:
        height = Mm(90)
        width = (height * img_demanda_descricao_width)/img_demanda_descricao_height
    else:
        width = Mm(140)
        height = (width * img_demanda_descricao_height)/img_demanda_descricao_width
    
    demanda_descricao_top = (prs.slide_height - height)/2
    demanda_descricao_left = (prs.slide_width - width)/2
    demanda_descricao_left = (prs.slide_width - width)/2
    demanda_descricao_slide = slide.shapes.add_picture(demanda_descricao_path, demanda_descricao_left, demanda_descricao_top, width=width)
    
    # Título - INFLUENCIADORES DE DEMANDA
    text_top = Mm(0)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'INFLUENCIADORES DE DEMANDA'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    if 'sim' in zoneamento_sim_nao:
        
        try:
            if zon_png_3:
                num_zon_png = 3
            elif zon_png_2:
                num_zon_png = 2
            elif zon_png_1:
                num_zon_png = 1


            for i in range(1, num_zon_png + 1):
            
                ##### SLIDE Zoneamento PNG
                # background
                blank_slide_layout = prs.slide_layouts[6]
                prs.slide_height = round(prs.slide_width * (9/16))
                slide = prs.slides.add_slide(blank_slide_layout)
                bg_left = bg_top = Inches(0)
                bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
                # coloca a imagem como background
                slide.shapes._spTree.remove(bg._element)
                slide.shapes._spTree.insert(2, bg._element)

                # Título - Zoneamento 
                text_top = Mm(-7)
                text_left = Mm(0)
                titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
                tf = titulo.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = 'ZONEAMENTO'
                p.alignment = PP_ALIGN.CENTER
                font = run.font
                font.name = 'Century Gothic'
                font.bold = True
                font.size = Pt(30)
                font.color.rgb = RGBColor(25,37,54)

                # imagem inserida zon_png
                zon_path = folder + str(globals()["zon_png_"+str(i)])
                img_zon = Image.open(zon_path)
                img_zon_width, img_zon_height = img_zon.size
                if img_zon_height > 288:
                    height = Mm(100)
                    width = (height * img_zon_width)/img_zon_height
                else:
                    width = Mm(180)
                    height = (width * img_zon_height)/img_zon_width
                zon_top = (prs.slide_height - height)/2 - Mm(5)
                zon_left = (prs.slide_width - width)/2
                zon_left = (prs.slide_width - width)/2
                zon = slide.shapes.add_picture(zon_path, zon_left, zon_top, width=width)

                # Numeracao da página
                text_top = Mm(127)
                text_left = Mm(0)
                rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
                tf = rodape.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = str(num_pag)
                p.alignment = PP_ALIGN.RIGHT
                font = run.font
                font.name = 'Century Gothic'
                font.size = Pt(12)
                font.color.rgb = RGBColor(255,255,255)
                num_pag += 1

                
        except:
            pass
        


    

###################################################################
    ##### SLIDE 13 - ZONEAMENTO MODALIDADES
    # background
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_left = bg_top = Inches(0)
    bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # Título - Zoneamento
    text_top = Mm(0)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'ZONEAMENTO'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)

    ##### verifica a quantidade de zoneamentos
    if 'sim' in zoneamento_sim_nao or 'Sim' in zoneamento_sim_nao  or 'SIM' in zoneamento_sim_nao:
        try:
            if zon_7:
                num_zon = 7
                string = zon_1 + '\n' + zon_2 + '\n' + zon_3 + '\n' + zon_4 + '\n' + zon_5 + '\n' + zon_6 + '\n' + zon_7
                top = 10
            elif zon_6:
                num_zon = 6
                string = zon_1 + '\n' + zon_2 + '\n' + zon_3 + '\n' + zon_4 + '\n' + zon_5 + '\n' + zon_6
                top = 20
            elif zon_5:
                num_zon = 5
                string = zon_1 + '\n' + zon_2 + '\n' + zon_3 + '\n' + zon_4 + '\n' + zon_5
                top = 30
            elif zon_4:
                num_zon = 4
                string = zon_1 + '\n' + zon_2 + '\n' + zon_3 + '\n' + zon_4
                top = 40
            elif zon_3:
                num_zon = 3
                string = zon_1 + '\n' + zon_2 + '\n' + zon_3
                top = 50
            elif zon_2:
                num_zon = 2
                string = zon_1 + '\n' + zon_2
                top = 60
            elif zon_1:
                num_zon = 1
                string = zon_1
                top = 70
            
            inc_top = 20
        
            ### Texto Zon_1
            text_top = Mm(top)
            text_left = Mm(20)
            texto = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width - Mm(40), height=prs.slide_height/2)
            tf = texto.text_frame
            tf.word_wrap = True
            p = tf.add_paragraph()
            p.width = prs.slide_width - Mm(40)
            run = p.add_run()
            run.text = string
            p.alignment = PP_ALIGN.LEFT
            font = run.font
            font.name = 'Century Gothic'
            font.size = Pt(24)
            font.color.rgb = RGBColor(25,37,54)
                
        except:
            print('nao existe')
            pass

        # Numeracao da página
        text_top = Mm(127)
        text_left = Mm(0)
        rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
        tf = rodape.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = str(num_pag)
        p.alignment = PP_ALIGN.RIGHT
        font = run.font
        font.name = 'Century Gothic'
        font.size = Pt(12)
        font.color.rgb = RGBColor(255,255,255)
        num_pag += 1

###################################################################
    ##### SLIDE 14 - RATING REGIÃO
    # background
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_left = bg_top = Inches(0)
    bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)


    # Título - Rating e Região
    text_top = Mm(0)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'RATING REGIÃO'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)

    # subtítulo - Rating Eemovel
    text_top = Mm(16)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'Rating EEmovel: ' + str(rating)
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(24)
    font.color.rgb = RGBColor(0,0,0)

    # imagem inserida 
    img_casa_path = path_slides_padrao + img_casa
    img_img_casa = Image.open(img_casa_path)
    img_img_casa_width, img_img_casa_height = img_img_casa.size
    img_casa_top = Mm(60)
    height = Mm(45)
    width = (height * img_img_casa_width)/img_img_casa_height
    img_casa_left = Mm(10)
    img_casa = slide.shapes.add_picture(img_casa_path, img_casa_left, img_casa_top, width=width)            


    # Texto Rating 
    text_top = Mm(30)
    text_left = Mm(70)
    texto = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width - Mm(75), height=prs.slide_height/2)
    tf = texto.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'É em função de fatores socioeconômicos, demográficos e de mercado imobiliário, como por exemplo: renda  domiciliar renda per capta, taxa alfabetização, longevidade,valor mediano imóveis absoluto e relativo. \n\nClassificação: A+++,A++...: É a Mediana da região mais o desvio padrão da região. Normalmente o Valor C é a mediana. A partir disto somamos o desvio padrão para cada classificação, exemplo: C = Mediana da região. C+ = Mediana + 1 desvio padrão. B= Mediana + 2 desvios padrão. B+ = Mediana + 3 desvios padrão. Quando o contrário, ou seja, D+: É a mediana - meio desvio padrão, e assim por diante.\n\nA estimativa populacional é realizada com o mecanismo de progressão geométrica. Mantendo tudo o mais constante (ceteris paribus), estimamos a taxa de crescimento populacional e colocamos na amostragem. Somente o crescimento é levado em conta neste caso, não consideramos migrações, mortalidade e natalidade no cálculo.'
    p.alignment = PP_ALIGN.LEFT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(25,37,54)

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    ##### SLIDE Análise Influência Residencial
    # background
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_left = bg_top = Mm(0)
    bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # Título
    text_top = Mm(-4)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'ANÁLISE DE INFLUÊNCIA - RESIDENCIAL'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)

    # imagem inserida influência residencial
    am_residencial_path = folder + str(zona_de_influencia_residencial)
    img_am_residencial = Image.open(am_residencial_path)
    img_am_residencial_width, img_am_residencial_height = img_am_residencial.size
    am_residencial_top = Mm(20)
    if img_am_residencial_height > 288:
        height = Mm(100)
        width = (height * img_am_residencial_width)/img_am_residencial_height
    else:
        width = Mm(180)
        height = (width * img_am_residencial_height)/img_am_residencial_width
    am_residencial_left = (prs.slide_width - width)/2
    am_residencial = slide.shapes.add_picture(am_residencial_path, am_residencial_left, am_residencial_top, width=width)


    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    ##### SLIDE Análise Entorno - Residencial
    # background
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_left = bg_top = Inches(0)
    bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # imagem inserida tabela residencial
    tabela_residencial_path = folder + tabela_residencial_png
    img_tabela_residencial = Image.open(tabela_residencial_path)
    img_tabela_residencial_width, img_tabela_residencial_height = img_tabela_residencial.size
    if img_tabela_residencial_height > 288:
        height = Mm(100)
        width = (height * img_tabela_residencial_width)/img_tabela_residencial_height
    else:
        width = Mm(180)
        height = (width * img_tabela_residencial_height)/img_tabela_residencial_width
    tabela_residencial_top = (prs.slide_height - height)/2
    tabela_residencial_left = (prs.slide_width - width)/2
    tabela_residencial_left = (prs.slide_width - width)/2
    tabela_residencial = slide.shapes.add_picture(tabela_residencial_path, tabela_residencial_left, tabela_residencial_top, width=width)

    # Título - Análise Entorno 
    text_top = Mm(-4)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'ANÁLISE ENTORNO'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(26)
    font.color.rgb = RGBColor(25,37,54)

    # Subtítulo - Apartamento Venda
    text_top = Mm(6)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'Apartamento Venda'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(16)
    font.color.rgb = RGBColor(255,0,21)

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    ##### SLIDE Análise Comercial
    # background
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_left = bg_top = Inches(0)
    bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # Título
    text_top = Mm(0)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'ANÁLISE DE INFLUÊNCIA - COMERCIAL'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)

    # imagem inserida Amostra comercial
    am_comercial_path = folder + zona_de_influencia_comercial
    img_am_comercial = Image.open(am_comercial_path)
    img_am_comercial_width, img_am_comercial_height = img_am_comercial.size
    am_comercial_top = Mm(20)
    if img_am_comercial_height > 288:
        height = Mm(100)
        width = (height * img_am_comercial_width)/img_am_comercial_height
    else:
        width = Mm(180)
        height = (width * img_am_comercial_height)/img_am_comercial_width
    am_comercial_left = (prs.slide_width - width)/2
    am_comercial = slide.shapes.add_picture(am_comercial_path, am_comercial_left, am_comercial_top, width=width)


    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    ##### SLIDE Análise Entorno Comercial Venda
    # background
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_left = bg_top = Inches(0)
    bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # Título - Análise Entorno 
    text_top = Mm(0)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'ANÁLISE ENTORNO'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)

    # Subtítulo - Comercial Venda
    text_top = Mm(15)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'Comercial Venda'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(18)
    font.color.rgb = RGBColor(255,0,21)

    # imagem inserida tabela comercial
    tabela_comercial_path = folder + tabela_comercial_png
    img_tabela_comercial = Image.open(tabela_comercial_path)
    img_tabela_comercial_width, img_tabela_comercial_height = img_tabela_comercial.size
    width = Mm(140)
    height = (width * img_tabela_comercial_height)/img_tabela_comercial_width
    tabela_comercial_top = (prs.slide_height - height)/2
    tabela_comercial_left = (prs.slide_width - width)/2
    tabela_comercial = slide.shapes.add_picture(tabela_comercial_path, tabela_comercial_left, tabela_comercial_top, width=width)

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1


###################################################################
    ##### SLIDE Análise de Influência Terreno
    # background
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_left = bg_top = Inches(0)
    bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # Título - Análise de Influência Terreno
    text_top = Mm(0)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'ANÁLISE DE INFLUÊNCIA TERRENO'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)


    # imagem inserida zona_de_influencia_terreno
    zona_de_influencia_terreno_path = folder + zona_de_influencia_terreno
    img_zona_de_influencia_terreno = Image.open(zona_de_influencia_terreno_path)
    img_zona_de_influencia_terreno_width, img_zona_de_influencia_terreno_height = img_zona_de_influencia_terreno.size
    width = Mm(130)
    height = (width * img_zona_de_influencia_terreno_height)/img_zona_de_influencia_terreno_width
    zona_de_influencia_terreno_top = (prs.slide_height - height)/2
    zona_de_influencia_terreno_left = (prs.slide_width - width)/2
    zona_de_influencia_terreno = slide.shapes.add_picture(zona_de_influencia_terreno_path, zona_de_influencia_terreno_left, zona_de_influencia_terreno_top, width=width)

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    ##### SLIDE Análise Entorno Terreno Venda
    # background
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_left = bg_top = Inches(0)
    bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # Título - Análise Entorno 
    text_top = Mm(0)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'ANÁLISE ENTORNO'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)

    # Subtítulo - Terreno Vendas
    text_top = Mm(15)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'Terreno Vendas'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(18)
    font.color.rgb = RGBColor(255,0,21)

    # imagem inserida tabela terreno
    tabela_terreno_path = folder + tabela_terreno_png
    img_tabela_terreno = Image.open(tabela_terreno_path)
    img_tabela_terreno_width, img_tabela_terreno_height = img_tabela_terreno.size
    width = Mm(140)
    height = (width * img_tabela_terreno_height)/img_tabela_terreno_width
    tabela_terreno_top = (prs.slide_height - height)/2
    tabela_terreno_left = (prs.slide_width - width)/2
    tabela_terreno = slide.shapes.add_picture(tabela_terreno_path, tabela_terreno_left, tabela_terreno_top, width=width)

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1


###################################################################
    ##### SLIDE Avaliação Evolutivo - Parte 1
    # background
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_left = bg_top = Inches(0)
    bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # Título - Avaliação Evolutivo - Parte 1
    text_top = Mm(0)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'AVALIAÇÃO EVOLUTIVO - PARTE 1'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)

    # Texto - Comercial Venda
    text_top = Mm(30)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'Foram levantados terrenos em oferta na mesma microrregião\n\nMEDIANA DOS TERRENOS OFERTADOS\nDENTRO DO RAIO DE INFLUÊNCIA'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(18)
    font.color.rgb = RGBColor(25,37,54)

    # Subtítulo - Comercial Venda
    text_top = Mm(80)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    b = valor_terreno_evolutivo
    a = locale.format_string("%1.2f", float(b), grouping=True)
    try:
        run.text = 'R$ ' + str(formated_number(mediana_terreno)) + ' x ' + str(formated_number(area_terreno)) + 'm² = R$ ' + str(a)
    except Exception as e:
        print(agencia,'Não encontrou número mediana_terreno ou area_terreno', e)
        pass
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(24)
    font.bold = True
    font.color.rgb = RGBColor(25,37,54)


    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    ##### SLIDE Avaliação Evolutivo - Parte 2 - Texto Inicial
    # background
    texto_avaliacao_evolutivo_parte2_path = path_slides_padrao + texto_avaliacao_evolutivo_parte_2
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    texto_avaliacao_evolutivo_parte2_left = texto_avaliacao_evolutivo_parte2_top = Inches(0)
    texto_avaliacao_evolutivo_parte2 = slide.shapes.add_picture(texto_avaliacao_evolutivo_parte2_path, texto_avaliacao_evolutivo_parte2_left, texto_avaliacao_evolutivo_parte2_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(texto_avaliacao_evolutivo_parte2._element)
    slide.shapes._spTree.insert(2, texto_avaliacao_evolutivo_parte2._element)

    # # Título - Avaliação Evolutivo - Parte 2
    # text_top = Mm(0)
    # text_left = Mm(0)
    # titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    # tf = titulo.text_frame
    # tf.word_wrap = True
    # p = tf.add_paragraph()
    # run = p.add_run()
    # run.text = 'AVALIAÇÃO EVOLUTIVO - PARTE 2'
    # p.alignment = PP_ALIGN.CENTER
    # font = run.font
    # font.name = 'Century Gothic'
    # font.bold = True
    # font.size = Pt(30)
    # font.color.rgb = RGBColor(25,37,54)

    

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    if avaliacao_evolutivo_parte_2:
    ###################################################################
        ##### SLIDE Avaliação Evolutivo - Parte 2 - Imagem
        # background
        blank_slide_layout = prs.slide_layouts[6]
        prs.slide_height = round(prs.slide_width * (9/16))
        slide = prs.slides.add_slide(blank_slide_layout)
        bg_left = bg_top = Inches(0)
        bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
        # coloca a imagem como background
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)

        # Título - Avaliação Evolutivo - Parte 2
        text_top = Mm(0)
        text_left = Mm(0)
        titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
        tf = titulo.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = 'AVALIAÇÃO EVOLUTIVO - PARTE 2'
        p.alignment = PP_ALIGN.CENTER
        font = run.font
        font.name = 'Century Gothic'
        font.bold = True
        font.size = Pt(30)
        font.color.rgb = RGBColor(25,37,54)

        # imagem inserida Amostra comercial
        avaliacao_evolutivo_parte2_path = folder + avaliacao_evolutivo_parte_2
        img_avaliacao_evolutivo_parte2 = Image.open(avaliacao_evolutivo_parte2_path)
        img_avaliacao_evolutivo_parte2_width, img_avaliacao_evolutivo_parte2_height = img_avaliacao_evolutivo_parte2.size
        avaliacao_evolutivo_parte2_top = Mm(25)
        height = Mm(90)
        width = (height * img_avaliacao_evolutivo_parte2_width)/img_avaliacao_evolutivo_parte2_height
        avaliacao_evolutivo_parte2_left = (prs.slide_width - width)/2
        avaliacao_evolutivo_parte2 = slide.shapes.add_picture(avaliacao_evolutivo_parte2_path, avaliacao_evolutivo_parte2_left, avaliacao_evolutivo_parte2_top, width=width)

        # Numeracao da página
        text_top = Mm(127)
        text_left = Mm(0)
        rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
        tf = rodape.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = str(num_pag)
        p.alignment = PP_ALIGN.RIGHT
        font = run.font
        font.name = 'Century Gothic'
        font.size = Pt(12)
        font.color.rgb = RGBColor(255,255,255)
        num_pag += 1

###################################################################
    ##### SLIDE Avaliação Evolutivo - Parte 2 - Cálculos
    # background
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_left = bg_top = Inches(0)
    bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # Título - Avaliação Evolutivo - Parte 2
    text_top = Mm(10)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'AVALIAÇÃO EVOLUTIVO - PARTE 2'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)


    # Subtítulo - Avaliação Evolutivo - Parte 2
    text_top = Mm(40)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    b = m2_construcao
    a = locale.format_string("%1.2f", float(b), grouping=True)
    run.text = 'CUSTO TOTAL DE CONSTRUÇÃO'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(18)
    font.color.rgb = RGBColor(25,37,54)


    # Subtítulo - Avaliação Evolutivo - Parte 2
    text_top = Mm(60)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    b = m2_construcao
    a = locale.format_string("%1.2f", float(b), grouping=True)
    run.text = 'PINI - Prédio sem elevador médio em ' + str(cidade_uf_referencia) + ' - R$ ' + str(a)
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(16)
    font.color.rgb = RGBColor(25,37,54)


    # Subtítulo - Avaliação Evolutivo - Parte 2
    text_top = Mm(85)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    a = m2_construcao
    formated_m2_construcao = locale.format_string("%1.2f", float(a), grouping=True)
    b = valor_construcao
    formated_valor_construcao = locale.format_string("%1.2f", float(b), grouping=True)
    c = area_construida
    formated_area_construida =locale.format_string("%1.2f", float(c), grouping=True)
    run.text = 'R$ ' + str(formated_m2_construcao) + ' x ' + str(formated_area_construida) + 'm² = R$ ' + str(formated_valor_construcao)
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(22)
    font.bold = True
    font.color.rgb = RGBColor(25,37,54)
    

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    ##### SLIDE Avaliação Evolutivo - Final
    # background
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_left = bg_top = Inches(0)
    bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # Título - Avaliação Evolutivo - Parte 2
    text_top = Mm(0)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'AVALIAÇÃO EVOLUTIVO - FINAL'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)


    # Texto - Avaliação Evolutivo Final - Terreno
    text_top = Mm(25)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    a = mediana_terreno
    formated_mediana_terreno = locale.format_string("%1.2f", float(a), grouping=True)
    b = area_terreno
    formated_area_terreno = locale.format_string("%1.2f", float(b), grouping=True)
    c = valor_terreno_evolutivo
    formated_valor_terreno_evolutivo = locale.format_string("%1.2f", float(c), grouping=True)
    run.text = 'Terreno\nR$ ' + str(formated_mediana_terreno) + ' x ' + str(formated_area_terreno) + 'm² = R$ ' + str(formated_valor_terreno_evolutivo)
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = False
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)

    # Texto - Avaliação Evolutivo Final - Custo de Construção
    text_top = Mm(45)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    a = m2_construcao
    formated_m2_construcao = locale.format_string("%1.2f", float(a), grouping=True)
    b = area_construida
    formated_area_construida = locale.format_string("%1.2f", float(b), grouping=True)
    c = valor_construcao
    formated_valor_construcao = locale.format_string("%1.2f", float(c), grouping=True)
    run.text = '\nCusto de Construção\nR$ ' + str(formated_m2_construcao) + ' x ' + str(formated_area_construida) + 'm² = R$ ' + str(formated_valor_construcao) 
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = False
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)

    # Texto - Avaliação Evolutivo Final - Valor Total
    text_top = Mm(95)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    a = valor_total_imovel
    formated_valor_total_imovel = locale.format_string("%1.2f", float(a), grouping=True)
    run.text = 'Valor Total: R$ ' + str(formated_valor_total_imovel)
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)
    

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    ##### SLIDE Valor Venal
    # background
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_left = bg_top = Inches(0)
    bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # Título - Avaliação Evolutivo - Parte 2
    text_top = Mm(0)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'VALOR VENAL'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)

    # imagem inserida Amostra comercial
    valor_venal_path = folder + valor_venal_png
    img_valor_venal = Image.open(valor_venal_path)
    img_valor_venal_width, img_valor_venal_height = img_valor_venal.size
    width = Mm(140)
    height = (width * img_valor_venal_height)/img_valor_venal_width
    valor_venal_top = (prs.slide_height - height)/2
    valor_venal_left = (prs.slide_width - width)/2
    valor_venal_left = (prs.slide_width - width)/2
    valor_venal_slide = slide.shapes.add_picture(valor_venal_path, valor_venal_left, valor_venal_top, width=width)

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    ##### SLIDE Avaliação Comparativo Fórmulas
    # background
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_left = bg_top = Inches(0)
    bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # Título - Avaliação Comparativo
    text_top = Mm(0)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'AVALIAÇÃO COMPARATIVO'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)

    # imagem inserida Avaliacao formulas 1
    avaliacao_comparativo_formulas_1_path = path_slides_padrao + avaliacao_comparativo_formulas_1
    img_avaliacao_comparativo_formulas_1 = Image.open(avaliacao_comparativo_formulas_1_path)
    img_avaliacao_comparativo_formulas_1_width, img_avaliacao_comparativo_formulas_1_height = img_avaliacao_comparativo_formulas_1.size
    avaliacao_comparativo_formulas_1_top = Mm(35)
    height = Mm(20)
    width = (height * img_avaliacao_comparativo_formulas_1_width)/img_avaliacao_comparativo_formulas_1_height
    avaliacao_comparativo_formulas_1_left = Mm(10)
    avaliacao_comparativo_formulas_1 = slide.shapes.add_picture(avaliacao_comparativo_formulas_1_path, avaliacao_comparativo_formulas_1_left, avaliacao_comparativo_formulas_1_top, width=width)
    
    # imagem inserida Avaliacao formulas 2
    avaliacao_comparativo_formulas_2_path = path_slides_padrao + avaliacao_comparativo_formulas_2
    img_avaliacao_comparativo_formulas_2 = Image.open(avaliacao_comparativo_formulas_2_path)
    img_avaliacao_comparativo_formulas_2_width, img_avaliacao_comparativo_formulas_2_height = img_avaliacao_comparativo_formulas_2.size
    avaliacao_comparativo_formulas_2_top = Mm(60)
    height = Mm(10)
    width = (height * img_avaliacao_comparativo_formulas_2_width)/img_avaliacao_comparativo_formulas_2_height
    avaliacao_comparativo_formulas_2_left = Mm(10)
    avaliacao_comparativo_formulas_2 = slide.shapes.add_picture(avaliacao_comparativo_formulas_2_path, avaliacao_comparativo_formulas_2_left, avaliacao_comparativo_formulas_2_top, width=width)

    # Texto - ¹n
    text_top = Mm(65)
    text_left = Mm(10)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=Mm(140), height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = '¹n = ' + str(n)
    p.alignment = PP_ALIGN.LEFT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(18)
    font.color.rgb = RGBColor(25,37,54)

    # Texto - Avaliação Comparativo
    text_top = Mm(75)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    a = m2comercial
    formated_m2comercial = locale.format_string("%1.2f", float(a), grouping=True)
    b = area_construida
    formated_area_construida = locale.format_string("%1.2f", float(b), grouping=True)
    try:
        run.text = 'R$ ' + str(formated_m2comercial) +'/m² x ' + str(formated_area_construida) + 'm² = R$ ' + str(formated_number(valor_comparativo))
    except Exception as e:
        print(agencia,'Não encontrou número m2comercial, ou area_contruida ou valor_comparativo', e)
        pass
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(24)
    font.bold = True
    font.color.rgb = RGBColor(25,37,54)

    # Texto - ¹foram encontradas 12 amostras de  anúncios de lojas 
    text_top = Mm(100)
    text_left = Mm(10)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=Mm(70), height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = '¹foram encontradas ' + str(n) + ' amostras de  anúncios de lojas'
    p.alignment = PP_ALIGN.LEFT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(10)
    font.color.rgb = RGBColor(25,37,54)

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    ##### SLIDE Avaliação Comparativo - Parececer Digital no Método Comparativo
    # background
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_left = bg_top = Inches(0)
    bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # Título - AVALIAÇÃO COMPARATIVO
    text_top = Mm(20)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'AVALIAÇÃO COMPARATIVO'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)


    # Texto - PARECER DIGITAL NO MÉTODO COMPARATIVO
    text_top = Mm(45)
    text_left = Mm(55)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=Mm(140), height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'PARECER DIGITAL NO MÉTODO COMPARATIVO'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(18)
    font.color.rgb = RGBColor(25,37,54)

    # Texto - Números
    text_top = Mm(70)
    text_left = Mm(55)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=Mm(140), height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    try:
        run.text = 'R$ ' + str(formated_number(valor_comparativo))
    except Exception as e:
        print(agencia,'Não encontrou número R$ valor_comparativo', e)
        pass
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    if involutivo_valor > 0:
        ##### SLIDE Estudo Involutivo Itaú
        # background
        blank_slide_layout = prs.slide_layouts[6]
        prs.slide_height = round(prs.slide_width * (9/16))
        slide = prs.slides.add_slide(blank_slide_layout)
        bg_left = bg_top = Inches(0)
        bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
        # coloca a imagem como background
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)

        # Título - Estudo Involutivo Itaú
        text_top = Mm(-4)
        text_left = Mm(0)
        titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
        tf = titulo.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = 'ESTUDO INVOLUTIVO ITAÚ'
        p.alignment = PP_ALIGN.CENTER
        font = run.font
        font.name = 'Century Gothic'
        font.bold = True
        font.size = Pt(30)
        font.color.rgb = RGBColor(25,37,54)

        # Subtítulo - Método Involutivo Express
        text_top = Mm(10)
        text_left = Mm(0)
        titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
        tf = titulo.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = 'Método Involutivo Express Itaú – Coeficiente máximo (Amostras Residenciais)'
        p.alignment = PP_ALIGN.CENTER
        font = run.font
        font.name = 'Century Gothic'
        font.size = Pt(18)
        font.color.rgb = RGBColor(25,37,54)

        # imagem inserida Involutivo
        involutivo_path = folder + involutivo_png
        img_involutivo = Image.open(involutivo_path)
        img_involutivo_width, img_involutivo_height = img_involutivo.size
        involutivo_top = Mm(40)
        width = Mm(200)
        height = (width * img_involutivo_height)/img_involutivo_width
        involutivo_left = (prs.slide_width - width)/2
        involutivo = slide.shapes.add_picture(involutivo_path, involutivo_left, involutivo_top, width=width)

        # # Texto Rodapé
        # text_top = Mm(95)
        # text_left = Mm(0)
        # titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
        # tf = titulo.text_frame
        # tf.word_wrap = True
        # p = tf.add_paragraph()
        # run = p.add_run()
        # a = involutivo_valor
        # formated_involutivo_valor = locale.format_string("%1.2f", float(a), grouping=True)
        # run.text = 'Valor aproximado do Terreno: R$ ' + str(formated_involutivo_valor)
        # p.alignment = PP_ALIGN.CENTER
        # font = run.font
        # font.name = 'Century Gothic'
        # font.size = Pt(22)
        # font.bold = True
        # font.color.rgb = RGBColor(25,37,54)

        # Numeracao da página
        text_top = Mm(127)
        text_left = Mm(0)
        rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
        tf = rodape.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = str(num_pag)
        p.alignment = PP_ALIGN.RIGHT
        font = run.font
        font.name = 'Century Gothic'
        font.size = Pt(12)
        font.color.rgb = RGBColor(255,255,255)
        num_pag += 1

###################################################################
        ##### SLIDE Estudo – Viabilidade de Incorporação/Vocação Imobiliária EEmovel Comercial
        # background
        blank_slide_layout = prs.slide_layouts[6]
        prs.slide_height = round(prs.slide_width * (9/16))
        slide = prs.slides.add_slide(blank_slide_layout)
        bg_left = bg_top = Inches(0)
        bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
        # coloca a imagem como background
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)

        # Título - Viabilidade de Incorporação/Vocação Imobiliária EEmovel Residencial
        text_top = Mm(0)
        text_left = Mm(0)
        titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
        tf = titulo.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = 'ESTUDO - VIABILIDADE DE\nINCORPORAÇÃO/VOCAÇÃO IMOBILIÁRIA\nEEMOVEL RESIDENCIAL'
        p.alignment = PP_ALIGN.CENTER
        font = run.font
        font.name = 'Century Gothic'
        font.bold = True
        font.size = Pt(30)
        font.color.rgb = RGBColor(25,37,54)


        # Texto Rodapé - Viabilidade de Incorporação/Vocação Imobiliária EEmovel Residencial
        text_top = Mm(43)
        text_left = Mm(0)
        titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
        tf = titulo.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = 'De acordo com zoneamento local e amostras de dados de\nanúncio da região o endereço comporta o seguinte\nprojeto hipotético '
        p.alignment = PP_ALIGN.CENTER
        font = run.font
        font.name = 'Century Gothic'
        font.size = Pt(16)
        font.color.rgb = RGBColor(25,37,54)


        # Texto - Projeto Hipotético: Sala Comercial
        text_top = Mm(70)
        text_left = Mm(10)
        titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
        tf = titulo.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = 'Projeto Hipotético: ' + proj_hiptotetico
        p.alignment = PP_ALIGN.LEFT
        font = run.font
        font.name = 'Century Gothic'
        font.size = Pt(18)
        font.color.rgb = RGBColor(25,37,54)

        # Texto - Área Útil
        text_top = Mm(80)
        text_left = Mm(10)
        titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
        tf = titulo.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        run = p.add_run()
        try:
            run.text = 'Área Útil: ' + str(formated_number(proj_area_util)) + ' m²'
        except Exception as e:
            # print(agencia,'Não encontrou número Área Útil: proj_area_util', e)
            pass
        p.alignment = PP_ALIGN.LEFT
        font = run.font
        font.name = 'Century Gothic'
        font.size = Pt(18)
        font.color.rgb = RGBColor(25,37,54)

        # Texto - Área Total
        text_top = Mm(90)
        text_left = Mm(10)
        titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
        tf = titulo.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        run = p.add_run()
        try:
            run.text = 'Área Total: ' + str(formated_number(proj_area_total)) + ' m²'
        except Exception as e:
            # print(agencia,'Não encontrou número Área Total: proj_area_total', e)
            pass
        p.alignment = PP_ALIGN.LEFT
        font = run.font
        font.name = 'Century Gothic'
        font.size = Pt(18)
        font.color.rgb = RGBColor(25,37,54)

        # Texto - Valor de Oferta
        text_top = Mm(100)
        text_left = Mm(10)
        titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
        tf = titulo.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        run = p.add_run()
        try:
            run.text = 'Valor de oferta: R$ ' + str(formated_number(valor_oferta))
        except Exception as e:
            # print(agencia,'Não encontrou número Valor de oferta', e)
            pass
        p.alignment = PP_ALIGN.LEFT
        font = run.font
        font.name = 'Century Gothic'
        font.size = Pt(18)
        font.color.rgb = RGBColor(25,37,54)

        

        # Numeracao da página
        text_top = Mm(127)
        text_left = Mm(0)
        rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
        tf = rodape.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = str(num_pag)
        p.alignment = PP_ALIGN.RIGHT
        font = run.font
        font.name = 'Century Gothic'
        font.size = Pt(12)
        font.color.rgb = RGBColor(255,255,255)
        num_pag += 1

###################################################################
        ##### SLIDE Cálculos Avaliatórios
        # background 
        calculos_avaliatorios_path = path_slides_padrao + calculos_avaliatorios
        blank_slide_layout = prs.slide_layouts[6]
        prs.slide_height = round(prs.slide_width * (9/16))
        slide = prs.slides.add_slide(blank_slide_layout)
        calculos_avaliatorios_left = calculos_avaliatorios_top = Inches(0)
        calculos_avaliatorios = slide.shapes.add_picture(calculos_avaliatorios_path, calculos_avaliatorios_left, calculos_avaliatorios_top, width=prs.slide_width, height=prs.slide_height)
        # coloca a imagem como background
        slide.shapes._spTree.remove(calculos_avaliatorios._element)
        slide.shapes._spTree.insert(2, calculos_avaliatorios._element)

        # Numeracao da página
        text_top = Mm(127)
        text_left = Mm(0)
        rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
        tf = rodape.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = str(num_pag)
        p.alignment = PP_ALIGN.RIGHT
        font = run.font
        font.name = 'Century Gothic'
        font.size = Pt(12)
        font.color.rgb = RGBColor(255,255,255)
        num_pag += 1        
    
    else:
        ##### SLIDE Estudo Involutivo EEmovel
        # background
        blank_slide_layout = prs.slide_layouts[6]
        prs.slide_height = round(prs.slide_width * (9/16))
        slide = prs.slides.add_slide(blank_slide_layout)
        bg_left = bg_top = Inches(0)
        bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
        # coloca a imagem como background
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)

        # Título - Estudo Involutivo EEmovel
        text_top = Mm(0)
        text_left = Mm(0)
        titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
        tf = titulo.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = 'ESTUDO INVOLUTIVO ITAÚ'
        p.alignment = PP_ALIGN.CENTER
        font = run.font
        font.name = 'Century Gothic'
        font.bold = True
        font.size = Pt(30)
        font.color.rgb = RGBColor(25,37,54)

        # Subtítulo - Método Involutivo Express Itaú
        text_top = Mm(17)
        text_left = Mm(0)
        titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
        tf = titulo.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = 'Método Involutivo Express Itaú – Coeficiente máximo (Amostras Residenciais)'
        p.alignment = PP_ALIGN.CENTER
        font = run.font
        font.name = 'Century Gothic'
        font.size = Pt(13)
        font.color.rgb = RGBColor(25,37,54)

        # imagem inserida Involutivo
        involutivo_path = folder + involutivo_png
        img_involutivo = Image.open(involutivo_path)
        img_involutivo_width, img_involutivo_height = img_involutivo.size
        involutivo_top = Mm(43)
        width = Mm(190)
        height = (width * img_involutivo_height)/img_involutivo_width
        involutivo_left = (prs.slide_width - width)/2
        involutivo = slide.shapes.add_picture(involutivo_path, involutivo_left, involutivo_top, width=width)

        
        # Numeracao da página
        text_top = Mm(127)
        text_left = Mm(0)
        rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
        tf = rodape.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = str(num_pag)
        p.alignment = PP_ALIGN.RIGHT
        font = run.font
        font.name = 'Century Gothic'
        font.size = Pt(12)
        font.color.rgb = RGBColor(255,255,255)
        num_pag += 1




###################################################################
    ##### SLIDE Considerações
    # background 
    consideracoes_path = path_slides_padrao + consideracoes
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    consideracoes_left = consideracoes_top = Inches(0)
    consideracoes = slide.shapes.add_picture(consideracoes_path, consideracoes_left, consideracoes_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(consideracoes._element)
    slide.shapes._spTree.insert(2, consideracoes._element)

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    ##### SLIDE Vocação Imobiliária
    # background
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_left = bg_top = Inches(0)
    bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)


    # Título - Vocação Imobiliária
    text_top = Mm(20)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'VOCAÇÃO IMOBILIÁRIA'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(30)
    font.color.rgb = RGBColor(25,37,54)

    # Texto - Vocação Imobiliária
    text_top = Mm(50)
    text_left = Mm(10)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width - Mm(20), height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = '\nApós levantamento de dados da região, os dados apontam que no endereço onde foram realizados os estudos o cenário que proporciona o melhor valor de venda é ' + str(parecer).lower() + '.'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(16)
    font.color.rgb = RGBColor(25,37,54)

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    ##### SLIDE Parceiros
    # background
    slide_5_path = path_slides_padrao + parceiros
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    slide_5_left = slide_5_top = Inches(0)
    slide_5 = slide.shapes.add_picture(slide_5_path, slide_5_left, slide_5_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(slide_5._element)
    slide.shapes._spTree.insert(2, slide_5._element)

    # Texto - Mais de 500 empresas do mercado imobiliário no Brasil
    text_top = Mm(90)
    text_left = Mm(140)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=Mm(80), height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'E mais de 500 empresas\ndo mercado imobiliário no Brasil'
    p.alignment = PP_ALIGN.LEFT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(16)
    font.color.rgb = RGBColor(25,37,54)

    # Numeracao da página
    text_top = Mm(127)
    text_left = Mm(0)
    rodape = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = rodape.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = str(num_pag)
    p.alignment = PP_ALIGN.RIGHT
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(12)
    font.color.rgb = RGBColor(255,255,255)
    num_pag += 1

###################################################################
    ##### SLIDE OBRIGADO
    # background
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height = round(prs.slide_width * (9/16))
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_left = bg_top = Inches(0)
    bg = slide.shapes.add_picture(bg_path, bg_left, bg_top, width=prs.slide_width, height=prs.slide_height)
    # coloca a imagem como background
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)


    # Título - OBRIGADO
    text_top = Mm(35)
    text_left = Mm(0)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width, height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'OBRIGADO'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.bold = True
    font.size = Pt(45)
    font.color.rgb = RGBColor(25,37,54)

    # Texto - Equipe Valuation Intelligence
    text_top = Mm(60)
    text_left = Mm(20)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width - Mm(40), height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'Equipe Valuation Intelligence'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(20)
    font.color.rgb = RGBColor(25,37,54)

    # Texto - Email
    text_top = Mm(75)
    text_left = Mm(20)
    titulo = slide.shapes.add_textbox(text_left, text_top, width=prs.slide_width - Mm(40), height=prs.slide_height/2)
    tf = titulo.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = 'valuation-intelligence@eemovel.com'
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.name = 'Century Gothic'
    font.size = Pt(16)
    font.color.rgb = RGBColor(25,37,54)
    printar = False
    if printar == True:
        print(type(agencia),'agencia',agencia)
        print(type(endereco),'endereco',endereco)
        print(type(foto_fachada_png),'foto_fachada_png',foto_fachada_png)
        print(type(area_terreno),'area_terreno',area_terreno)
        print(type(area_construida),'area_construida',area_construida)
        print(type(iptu_espelho),'iptu_espelho',iptu_espelho)
        print(type(mapa_demanda),'mapa_demanda',mapa_demanda)
        print(type(demanda_descricao),'demanda_descricao',demanda_descricao)
        print(type(mapa_satelite_png),'mapa_satelite_png',mapa_satelite_png)
        print(type(zoneamento_sim_nao),'zoneamento_sim_nao',zoneamento_sim_nao)
        print(type(zon_1),'zon_1',zon_1)
        print(type(zon_2),'zon_2',zon_2)
        print(type(zon_3),'zon_3',zon_3)
        print(type(zon_4),'zon_4',zon_4)
        print(type(zon_5),'zon_5',zon_5)
        print(type(zon_6),'zon_6',zon_6)
        print(type(zon_7),'zon_7',zon_7)
        print(type(zon_png_1),'zon_png_1',zon_png_1)
        print(type(zon_png_2),'zon_png_2',zon_png_2)
        print(type(zon_png_3),'zon_png_3',zon_png_3)
        print(type(rating),'rating',rating)
        print(type(zona_de_influencia_residencial),'zona_de_influencia_residencial',zona_de_influencia_residencial)
        print(type(tabela_residencial_png),'tabela_residencial_png',tabela_residencial_png)
        print(type(zona_de_influencia_comercial),'zona_de_influencia_comercial',zona_de_influencia_comercial)
        print(type(tabela_comercial_png),'tabela_comercial_png',tabela_comercial_png)
        print(type(zona_de_influencia_terreno),'zona_de_influencia_terreno',zona_de_influencia_terreno)
        print(type(tabela_terreno_png),'tabela_terreno_png',tabela_terreno_png)
        print(type(mediana_terreno),'mediana_terreno',mediana_terreno)
        print(type(valor_terreno_evolutivo),'valor_terreno_evolutivo',valor_terreno_evolutivo)
        print(type(avaliacao_evolutivo_parte_2),'avaliacao_evolutivo_parte_2',avaliacao_evolutivo_parte_2)
        print(type(cidade_uf_referencia),'cidade_uf_referencia',cidade_uf_referencia)
        print(type(m2_construcao),'m2_construcao',m2_construcao)
        print(type(valor_construcao),'valor_construcao',valor_construcao)
        print(type(valor_total_imovel),'valor_total_imovel',valor_total_imovel)
        print(type(valor_venal_png),'valor_venal_png',valor_venal_png)
        print(type(n),'n',n)
        print(type(m2comercial),'m2comercial',m2comercial)
        print(type(valor_comparativo),'valor_comparativo',valor_comparativo)
        print(type(involutivo_png),'involutivo_png',involutivo_png)
        print(type(involutivo_valor),'involutivo_valor',involutivo_valor)
        print(type(proj_hiptotetico),'proj_hiptotetico',proj_hiptotetico)
        print(type(proj_area_util),'proj_area_util',proj_area_util)
        print(type(proj_area_total),'proj_area_total',proj_area_total)
        print(type(valor_oferta),'valor_oferta',valor_oferta)
        print(type(parecer),'parecer',parecer)
        print(type(finalizado),'finalizado',finalizado)
    print('SUCESSO!!', agencia, cidade_uf_referencia)
    return(prs.save(agencia + '_apresentacao.pptx'))
    # return(lista)

if __name__ == "__main__":

    
    tabela_geral = xlrd.open_workbook('create_apresentacao_eemovel.xlsx')
    sheet = tabela_geral.sheet_by_name('Página1')
    num_rows = sheet.nrows
    # print(num_rows)
    lista = []
    lista_agencias = []
    folders = os.listdir('./')
    lista_folders = []
    for i in range(1,num_rows):
        # print('+++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++')
        try:
            agencia = sheet.cell_value(i,0)
            # print('agencia:', type(agencia), agencia)
        except Exception as e:
            # print('Não encontrou agencia', e)
            pass
        try:
            endereco = sheet.cell_value(i,1)
            # print('endereco:', type(endereco), endereco)
        except Exception as e:
            # print('Não encontrou endereco', e)
            pass
        try:
            foto_fachada_png = sheet.cell_value(i,2)
            # print('foto_fachada_png:', type(foto_fachada_png), foto_fachada_png)
        except Exception as e:
            # print('Não encontrou foto_fachada_png', e)
            pass
        try:
            area_terreno = sheet.cell_value(i,3)
            # print('area_terreno:', type(area_terreno), area_terreno)
        except Exception as e:
            # print('Não encontrou area_terreno', e)
            pass
        try:
            area_construida = sheet.cell_value(i,4)
            # print('area_construida:', type(area_construida), area_construida)
        except Exception as e:
            # print('Não encontrou area_construida', e)
            pass
        try:
            iptu_espelho = sheet.cell_value(i,5)
            # print('iptu_espelho:', type(iptu_espelho), iptu_espelho)
        except Exception as e:
            # print('Não encontrou iptu_espelho', e)
            pass
        try:
            mapa_demanda = sheet.cell_value(i,6)
            # print('mapa_demanda:', type(mapa_demanda), mapa_demanda)
        except Exception as e:
            # print('Não encontrou mapa_demanda', e)
            pass
        try:
            demanda_descricao = sheet.cell_value(i,7)
            # print('demanda_descricao:', type(demanda_descricao), demanda_descricao)
        except Exception as e:
            # print('Não encontrou demanda_descricao', e)
            pass
        try:
            mapa_satelite_png = sheet.cell_value(i,8)
            # print('mapa_satelite_png:', type(mapa_satelite_png), mapa_satelite_png)
        except Exception as e:
            # print('Não encontrou mapa_satelite_png', e)
            pass
        try:
            zoneamento_sim_nao = sheet.cell_value(i,9)
            # print('zoneamento_sim_nao:', type(zoneamento_sim_nao), zoneamento_sim_nao)
        except Exception as e:
            # print('Não encontrou zoneamento_sim_nao', e)
            pass
        try:
            zon_1 = sheet.cell_value(i,10)
            # print('zon_1:', type(zon_1), zon_1)
        except Exception as e:
            # print('Não encontrou zon_1', e)
            pass
        try:
            zon_2 = sheet.cell_value(i,11)
            # print('zon_2:', type(zon_2), zon_2)
        except Exception as e:
            # print('Não encontrou zon_2', e)
            pass
        try:
            zon_3 = sheet.cell_value(i,12)
            # print('zon_3:', type(zon_3), zon_3)
        except Exception as e:
            # print('Não encontrou zon_3', e)
            pass
        try:
            zon_4 = sheet.cell_value(i,13)
            # print('zon_4:', type(zon_4), zon_4)
        except Exception as e:
            # print('Não encontrou zon_4', e)
            pass
        try:
            zon_5 = sheet.cell_value(i,14)
            # print('zon_5:', type(zon_5), zon_5)
        except Exception as e:
            # print('Não encontrou zon_5', e)
            pass
        try:
            zon_6 = sheet.cell_value(i,15)
            # print('zon_6:', type(zon_6), zon_6)
        except Exception as e:
            # print('Não encontrou zon_6', e)
            pass
        try:
            zon_7 = sheet.cell_value(i,16)
            # print('zon_7:', type(zon_7), zon_7)
        except Exception as e:
            # print('Não encontrou zon_7', e)
            pass
        try:
            zon_png_1 = sheet.cell_value(i,17)
            # print('zon_png_1:', type(zon_png_1), zon_png_1)
        except Exception as e:
            # print('Não encontrou zon_png_1', e)
            pass
        try:
            zon_png_2 = sheet.cell_value(i,18)
            # print('zon_png_2:', type(zon_png_2), zon_png_2)
        except Exception as e:
            # print('Não encontrou zon_png_2', e)
            pass
        try:
            zon_png_3 = sheet.cell_value(i,19)
            # print('zon_png_3:', type(zon_png_3), zon_png_3)
        except Exception as e:
            # print('Não encontrou zon_png_3', e)
            pass
        try:
            rating = sheet.cell_value(i,20)
            # print('rating:', type(rating), rating)
        except Exception as e:
            # print('Não encontrou rating', e)
            pass
        try:
            zona_de_influencia_residencial = sheet.cell_value(i,21)
            # print('zona_de_influencia_residencial:', type(zona_de_influencia_residencial), zona_de_influencia_residencial)
        except Exception as e:
            # print('Não encontrou zona_de_influencia_residencial', e)
            pass
        try:
            tabela_residencial_png = sheet.cell_value(i,22)
            # print('tabela_residencial_png:', type(tabela_residencial_png), tabela_residencial_png)
        except Exception as e:
            # print('Não encontrou tabela_residencial_png', e)
            pass
        try:
            zona_de_influencia_comercial = sheet.cell_value(i,23)
            # print('zona_de_influencia_comercial:', type(zona_de_influencia_comercial), zona_de_influencia_comercial)
        except Exception as e:
            # print('Não encontrou zona_de_influencia_comercial', e)
            pass
        try:
            tabela_comercial_png = sheet.cell_value(i,24)
            # print('tabela_comercial_png:', type(tabela_comercial_png), tabela_comercial_png)
        except Exception as e:
            # print('Não encontrou tabela_comercial_png', e)
            pass
        try:
            zona_de_influencia_terreno = sheet.cell_value(i,25)
            # print('zona_de_influencia_terreno:', type(zona_de_influencia_terreno), zona_de_influencia_terreno)
        except Exception as e:
            # print('Não encontrou zona_de_influencia_terreno', e)
            pass
        try:
            tabela_terreno_png = sheet.cell_value(i,26)
            # print('tabela_terreno_png:', type(tabela_terreno_png), tabela_terreno_png)
        except Exception as e:
            # print('Não encontrou tabela_terreno_png', e)
            pass
        try:
            mediana_terreno = sheet.cell_value(i,27)
            # print('mediana_terreno:', type(mediana_terreno), mediana_terreno)
        except Exception as e:
            # print('Não encontrou mediana_terreno', e)
            pass
        try:
            valor_terreno_evolutivo = sheet.cell_value(i,28)
            # print('valor_terreno_evolutivo:', type(valor_terreno_evolutivo), valor_terreno_evolutivo)
        except Exception as e:
            # print('Não encontrou valor_terreno_evolutivo', e)
            pass
        try:
            avaliacao_evolutivo_parte_2 = sheet.cell_value(i,29)
            # print('avaliacao_evolutivo_parte_2:', type(avaliacao_evolutivo_parte_2), avaliacao_evolutivo_parte_2)
        except Exception as e:
            # print('Não encontrou avaliacao_evolutivo_parte_2', e)
            pass
        try:
            cidade_uf_referencia = sheet.cell_value(i,30)
            # print('cidade_uf_referencia:', type(cidade_uf_referencia), cidade_uf_referencia)
        except Exception as e:
            # print('Não encontrou cidade_uf_referencia', e)
            pass
        try:
            m2_construcao = sheet.cell_value(i,31)
            # print('m2_construcao:', type(m2_construcao), m2_construcao)
        except Exception as e:
            # print('Não encontrou m2_construcao', e)
            pass
        try:
            valor_construcao = sheet.cell_value(i,32)
            # print('valor_construcao:', type(valor_construcao), valor_construcao)
        except Exception as e:
            # print('Não encontrou valor_construcao', e)
            pass
        try:
            valor_total_imovel = sheet.cell_value(i,33)
            # print('valor_total_imovel:', type(valor_total_imovel), valor_total_imovel)
        except Exception as e:
            # print('Não encontrou valor_total_imovel', e)
            pass
        try:
            valor_venal_png = sheet.cell_value(i,34)
            # print('valor_venal_png:', type(valor_venal_png), valor_venal_png)
        except Exception as e:
            # print('Não encontrou valor_venal_png', e)
            pass
        try:
            n = sheet.cell_value(i,35)
            # print('n:', type(n), n)
        except Exception as e:
            # print('Não encontrou n', e)
            pass
        try:
            m2comercial = sheet.cell_value(i,36)
            # print('m2comercial:', type(m2comercial), m2comercial)
        except Exception as e:
            # print('Não encontrou m2comercial', e)
            pass
        try:
            valor_comparativo = sheet.cell_value(i,37)
            # print('valor_comparativo:', type(valor_comparativo), valor_comparativo)
        except Exception as e:
            # print('Não encontrou valor_comparativo', e)
            pass
        try:
            involutivo_png = sheet.cell_value(i,38)
            # print('involutivo_png:', type(involutivo_png), involutivo_png)
        except Exception as e:
            # print('Não encontrou involutivo_png', e)
            pass
        try:
            involutivo_valor = sheet.cell_value(i,39)
            # print('involutivo_valor:', type(involutivo_valor), involutivo_valor)
        except Exception as e:
            # print('Não encontrou involutivo_valor', e)
            pass
        try:
            proj_hiptotetico = sheet.cell_value(i,40)
            # print('proj_hiptotetico:', type(proj_hiptotetico), proj_hiptotetico)
        except Exception as e:
            # print('Não encontrou proj_hiptotetico', e)
            pass
        try:
            proj_area_util = sheet.cell_value(i,41)
            # print('proj_area_util:', type(proj_area_util), proj_area_util)
        except Exception as e:
            # print('Não encontrou proj_area_util', e)
            pass
        try:
            proj_area_total = sheet.cell_value(i,42)
            # print('proj_area_total:', type(proj_area_total), proj_area_total)
        except Exception as e:
            # print('Não encontrou proj_area_total', e)
            pass
        try:
            valor_oferta = sheet.cell_value(i,43)
            # print('valor_oferta:', type(valor_oferta), valor_oferta)
        except Exception as e:
            # print('Não encontrou valor_oferta', e)
            pass
        try:
            parecer = sheet.cell_value(i,44)
            # print('parecer:', type(parecer), parecer)
        except Exception as e:
            # print('Não encontrou parecer', e)
            pass
        try:
            finalizado = sheet.cell_value(i,45)
            # print('finalizado:', type(finalizado), finalizado)
            finalizado = str(finalizado).lower()
        except Exception as e:
            # print('Não encontrou finalizado', e)
            pass
        # print('+++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++')
    
        if 'sim' in finalizado or 'Sim' in finalizado:
            for folder in folders:
                if '.' not in folder and folder not in lista_folders and str(folder).startswith(str(agencia)) and len(agencia) > 0:
                    lista_folders.append(folder)
            if str(agencia) in str(lista_folders) and len(str(agencia)) > 0:
                # print('++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++')
                # print('++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++')
                # print('agencia:', type(agencia), agencia)
                # print('endereco:', type(endereco), endereco)
                # print('foto_fachada_png:', type(foto_fachada_png), foto_fachada_png)
                # print('area_terreno:', type(area_terreno), area_terreno)
                # print('area_construida:', type(area_construida), area_construida)
                # print('iptu_espelho:', type(iptu_espelho), iptu_espelho)
                # print('mapa_demanda:', type(mapa_demanda), mapa_demanda)
                # print('demanda_descricao:', type(demanda_descricao), demanda_descricao)
                # print('mapa_satelite_png:', type(mapa_satelite_png), mapa_satelite_png)
                # print('zoneamento_sim_nao:', type(zoneamento_sim_nao), zoneamento_sim_nao)
                # print('zon_1:', type(zon_1), zon_1)
                # print('zon_2:', type(zon_2), zon_2)
                # print('zon_3:', type(zon_3), zon_3)
                # print('zon_4:', type(zon_4), zon_4)
                # print('zon_5:', type(zon_5), zon_5)
                # print('zon_6:', type(zon_6), zon_6)
                # print('zon_7:', type(zon_7), zon_7)
                # print('zon_png_1:', type(zon_png_1), zon_png_1)
                # print('zon_png_2:', type(zon_png_2), zon_png_2)
                # print('zon_png_3:', type(zon_png_3), zon_png_3)
                # print('rating:', type(rating), rating)
                # print('zona_de_influencia_residencial:', type(zona_de_influencia_residencial), zona_de_influencia_residencial)
                # print('tabela_residencial_png:', type(tabela_residencial_png), tabela_residencial_png)
                # print('zona_de_influencia_comercial:', type(zona_de_influencia_comercial), zona_de_influencia_comercial)
                # print('tabela_comercial_png:', type(tabela_comercial_png), tabela_comercial_png)
                # print('zona_de_influencia_terreno:', type(zona_de_influencia_terreno), zona_de_influencia_terreno)
                # print('tabela_terreno_png:', type(tabela_terreno_png), tabela_terreno_png)
                # print('mediana_terreno:', type(mediana_terreno), mediana_terreno)
                # print('valor_terreno_evolutivo:', type(valor_terreno_evolutivo), valor_terreno_evolutivo)
                # print('avaliacao_evolutivo_parte_2:', type(avaliacao_evolutivo_parte_2), avaliacao_evolutivo_parte_2)
                # print('cidade_uf_referencia:', type(cidade_uf_referencia), cidade_uf_referencia)
                # print('m2_construcao:', type(m2_construcao), m2_construcao)
                # print('valor_construcao:', type(valor_construcao), valor_construcao)
                # print('valor_total_imovel:', type(valor_total_imovel), valor_total_imovel)
                # print('valor_venal_png:', type(valor_venal_png), valor_venal_png)
                # print('n:', type(n), n)
                # print('m2comercial:', type(m2comercial), m2comercial)
                # print('valor_comparativo:', type(valor_comparativo), valor_comparativo)
                # print('involutivo_png:', type(involutivo_png), involutivo_png)
                # print('involutivo_valor:', type(involutivo_valor), involutivo_valor)
                # print('proj_hiptotetico:', type(proj_hiptotetico), proj_hiptotetico)
                # print('proj_area_util:', type(proj_area_util), proj_area_util)
                # print('proj_area_total:', type(proj_area_total), proj_area_total)
                # print('valor_oferta:', type(valor_oferta), valor_oferta)
                # print('parecer:', type(parecer), parecer)
                # print('finalizado:', type(finalizado), finalizado)
                # print('++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++')
                # print('++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++')
                
                try:
                    apresentacao(agencia,endereco,foto_fachada_png,area_construida,area_terreno,iptu_espelho,mapa_demanda,demanda_descricao,mapa_satelite_png,zoneamento_sim_nao,zon_1,zon_2,zon_3,zon_4,zon_5,zon_6,zon_7,zon_png_1,zon_png_2,zon_png_3,rating,zona_de_influencia_residencial,tabela_residencial_png,zona_de_influencia_comercial,tabela_comercial_png,zona_de_influencia_terreno,tabela_terreno_png,mediana_terreno,valor_terreno_evolutivo,avaliacao_evolutivo_parte_2,cidade_uf_referencia,m2_construcao,valor_construcao,valor_total_imovel,valor_venal_png,n,m2comercial,valor_comparativo,involutivo_png,involutivo_valor,proj_hiptotetico,proj_area_util,proj_area_total,valor_oferta,parecer,finalizado)
                except Exception as e:
                    print('Não rodou a agência', agencia, cidade_uf_referencia, e)
                        
                
           
        else:
            pass
            # if len(agencia) > 0:
            #     print(agencia, 'Não encontrado coluna FINALIZADO')
                
         
        